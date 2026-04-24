from pathlib import Path

from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "deliverables"
OUT_DIR.mkdir(exist_ok=True)


DATA = {
    "Visa": {
        "revenue_2025": 40000000000,
        "net_income_2025": 20058000000,
        "operating_income_2025": 23994000000,
        "interest_expense_2025": 589000000,
        "assets_2024": 94511000000,
        "assets_2025": 99627000000,
        "current_assets_2025": 37766000000,
        "cash_2025": 17164000000,
        "ar_2024": 2561000000,
        "ar_2025": 3126000000,
        "current_liabilities_2025": 35048000000,
        "equity_2024": 39137000000,
        "equity_2025": 37909000000,
        "total_debt_2025": 25171000000,
        "lt_debt_2025": 19602000000,
        "ppe_2024": 3824000000,
        "ppe_2025": 4236000000,
        "goodwill_2025": 19879000000,
        "intangibles_2025": 27646000000,
        "dividends_2025": 4634000000,
        "eps_2025": 10.20,
        "market_price_2026_04_24": 308.88,
        "shares_outstanding_2025": 1687629770,
    },
    "Mastercard": {
        "revenue_2025": 32791000000,
        "net_income_2025": 14968000000,
        "operating_income_2025": 18897000000,
        "interest_expense_2025": 722000000,
        "assets_2024": 48081000000,
        "assets_2025": 54157000000,
        "current_assets_2025": 23558000000,
        "cash_2025": 10566000000,
        "ar_2024": 3773000000,
        "ar_2025": 4609000000,
        "current_liabilities_2025": 22762000000,
        "equity_2024": 6485000000,
        "equity_2025": 7737000000,
        "total_debt_2025": 19000000000,
        "lt_debt_2025": 18251000000,
        "ppe_2024": 2138000000,
        "ppe_2025": 2303000000,
        "goodwill_2025": 9560000000,
        "intangibles_2025": 5382000000,
        "dividends_2025": 2756000000,
        "eps_2025": 16.52,
        "market_price_2026_04_24": 502.38,
        "shares_outstanding_2025": 885216386,
    },
}


COMMENTS = {
    "Net PP&E": "Mastercard runs with a smaller fixed-asset base than Visa, consistent with a lighter operating footprint.",
    "Goodwill and Intangibles": "Visa carries a much larger intangible base, which matters when interpreting asset structure and turnover.",
    "Total Long-term Debt": "Visa has a higher absolute amount of long-term debt, but Mastercard is more leveraged relative to equity.",
    "Total Debt / Total Assets": "Lower is less risky. Mastercard carries more debt relative to assets.",
    "Long-term Debt / Equity": "Mastercard looks far more leveraged because its equity base is small after buybacks.",
    "Current Ratio": "Both are just above 1.0, with Visa slightly safer on this measure.",
    "Quick Ratio": "Both are below 1.0; Mastercard is stronger on immediate liquidity.",
    "Interest Coverage": "Both firms comfortably cover interest expense, but Visa has the stronger safety margin.",
    "Net Profit Margin": "Visa converts a larger share of revenue into profit.",
    "ROE": "Mastercard's ROE is boosted by a very small equity base, so read with caution.",
    "ROA": "Mastercard generates more profit per dollar of assets.",
    "Asset Turnover": "Mastercard uses its asset base more efficiently to produce revenue.",
    "Receivables Turnover": "Visa collects or turns receivables faster.",
    "Inventory Turnover": "n/a for both: payment networks are service businesses with no meaningful inventory.",
    "Fixed Assets Turnover": "Mastercard generates more revenue per dollar of fixed assets.",
    "P/E": "Very similar valuation multiples based on current market prices and FY2025 EPS.",
    "Price to Book": "Mastercard's price-to-book is inflated by its low accounting equity after repurchases.",
    "Dividend Payout": "Visa returns a higher share of profits as dividends.",
}


RATIOS = [
    ("Net PP&E", "magnitude", "=RawData!U{row}", "=RawData!AR{row}"),
    ("Goodwill and Intangibles", "magnitude", "=RawData!V{row}+RawData!W{row}", "=RawData!AS{row}+RawData!AT{row}"),
    ("Total Long-term Debt", "magnitude", "=RawData!P{row}", "=RawData!AM{row}"),
    ("Total Debt / Total Assets", "structure", "=RawData!C{row}/RawData!E{row}", "=RawData!Z{row}/RawData!AB{row}"),
    ("Long-term Debt / Equity", "structure", "=RawData!P{row}/RawData!I{row}", "=RawData!AM{row}/RawData!AF{row}"),
    ("Current Ratio", "liquidity", "=RawData!F{row}/RawData!G{row}", "=RawData!AC{row}/RawData!AD{row}"),
    ("Quick Ratio", "liquidity", "=(RawData!Q{row}+RawData!K{row})/RawData!G{row}", "=(RawData!AN{row}+RawData!AH{row})/RawData!AD{row}"),
    ("Interest Coverage", "liquidity", "=RawData!R{row}/RawData!S{row}", "=RawData!AO{row}/RawData!AP{row}"),
    ("Net Profit Margin", "profitability", "=RawData!B{row}/RawData!A{row}", "=RawData!Y{row}/RawData!X{row}"),
    ("ROE", "profitability", "=RawData!B{row}/AVERAGE(RawData!H{row},RawData!I{row})", "=RawData!Y{row}/AVERAGE(RawData!AE{row},RawData!AF{row})"),
    ("ROA", "profitability", "=RawData!B{row}/AVERAGE(RawData!D{row},RawData!E{row})", "=RawData!Y{row}/AVERAGE(RawData!AA{row},RawData!AB{row})"),
    ("Asset Turnover", "efficiency", "=RawData!A{row}/AVERAGE(RawData!D{row},RawData!E{row})", "=RawData!X{row}/AVERAGE(RawData!AA{row},RawData!AB{row})"),
    ("Receivables Turnover", "efficiency", "=RawData!A{row}/AVERAGE(RawData!J{row},RawData!K{row})", "=RawData!X{row}/AVERAGE(RawData!AG{row},RawData!AH{row})"),
    ("Inventory Turnover", "efficiency", None, None),
    ("Fixed Assets Turnover", "efficiency", "=RawData!A{row}/AVERAGE(RawData!T{row},RawData!U{row})", "=RawData!X{row}/AVERAGE(RawData!AQ{row},RawData!AR{row})"),
    ("P/E", "market", "=RawData!N{row}/RawData!M{row}", "=RawData!AK{row}/RawData!AJ{row}"),
    ("Price to Book", "market", "=RawData!N{row}/(RawData!I{row}/RawData!O{row})", "=RawData!AK{row}/(RawData!AF{row}/RawData!AL{row})"),
    ("Dividend Payout", "market", "=RawData!L{row}/RawData!B{row}", "=RawData!AI{row}/RawData!Y{row}"),
]


def build_workbook():
    wb = Workbook()
    ws = wb.active
    ws.title = "RawData"

    headers = [
        "Visa Revenue 2025",
        "Visa Net Income 2025",
        "Visa Total Debt 2025",
        "Visa Assets 2024",
        "Visa Assets 2025",
        "Visa Current Assets 2025",
        "Visa Current Liabilities 2025",
        "Visa Equity 2024",
        "Visa Equity 2025",
        "Visa AR 2024",
        "Visa AR 2025",
        "Visa Dividends 2025",
        "Visa EPS 2025",
        "Visa Market Price (2026-04-24)",
        "Visa Shares Outstanding 2025",
        "Visa LT Debt 2025",
        "Visa Cash 2025",
        "Visa Operating Income 2025",
        "Visa Interest Expense 2025",
        "Visa PPE 2024",
        "Visa PPE 2025",
        "Visa Goodwill 2025",
        "Visa Intangibles 2025",
        "Mastercard Revenue 2025",
        "Mastercard Net Income 2025",
        "Mastercard Total Debt 2025",
        "Mastercard Assets 2024",
        "Mastercard Assets 2025",
        "Mastercard Current Assets 2025",
        "Mastercard Current Liabilities 2025",
        "Mastercard Equity 2024",
        "Mastercard Equity 2025",
        "Mastercard AR 2024",
        "Mastercard AR 2025",
        "Mastercard Dividends 2025",
        "Mastercard EPS 2025",
        "Mastercard Market Price (2026-04-24)",
        "Mastercard Shares Outstanding 2025",
        "Mastercard LT Debt 2025",
        "Mastercard Cash 2025",
        "Mastercard Operating Income 2025",
        "Mastercard Interest Expense 2025",
        "Mastercard PPE 2024",
        "Mastercard PPE 2025",
        "Mastercard Goodwill 2025",
        "Mastercard Intangibles 2025",
    ]
    values = [
        DATA["Visa"]["revenue_2025"],
        DATA["Visa"]["net_income_2025"],
        DATA["Visa"]["total_debt_2025"],
        DATA["Visa"]["assets_2024"],
        DATA["Visa"]["assets_2025"],
        DATA["Visa"]["current_assets_2025"],
        DATA["Visa"]["current_liabilities_2025"],
        DATA["Visa"]["equity_2024"],
        DATA["Visa"]["equity_2025"],
        DATA["Visa"]["ar_2024"],
        DATA["Visa"]["ar_2025"],
        DATA["Visa"]["dividends_2025"],
        DATA["Visa"]["eps_2025"],
        DATA["Visa"]["market_price_2026_04_24"],
        DATA["Visa"]["shares_outstanding_2025"],
        DATA["Visa"]["lt_debt_2025"],
        DATA["Visa"]["cash_2025"],
        DATA["Visa"]["operating_income_2025"],
        DATA["Visa"]["interest_expense_2025"],
        DATA["Visa"]["ppe_2024"],
        DATA["Visa"]["ppe_2025"],
        DATA["Visa"]["goodwill_2025"],
        DATA["Visa"]["intangibles_2025"],
        DATA["Mastercard"]["revenue_2025"],
        DATA["Mastercard"]["net_income_2025"],
        DATA["Mastercard"]["total_debt_2025"],
        DATA["Mastercard"]["assets_2024"],
        DATA["Mastercard"]["assets_2025"],
        DATA["Mastercard"]["current_assets_2025"],
        DATA["Mastercard"]["current_liabilities_2025"],
        DATA["Mastercard"]["equity_2024"],
        DATA["Mastercard"]["equity_2025"],
        DATA["Mastercard"]["ar_2024"],
        DATA["Mastercard"]["ar_2025"],
        DATA["Mastercard"]["dividends_2025"],
        DATA["Mastercard"]["eps_2025"],
        DATA["Mastercard"]["market_price_2026_04_24"],
        DATA["Mastercard"]["shares_outstanding_2025"],
        DATA["Mastercard"]["lt_debt_2025"],
        DATA["Mastercard"]["cash_2025"],
        DATA["Mastercard"]["operating_income_2025"],
        DATA["Mastercard"]["interest_expense_2025"],
        DATA["Mastercard"]["ppe_2024"],
        DATA["Mastercard"]["ppe_2025"],
        DATA["Mastercard"]["goodwill_2025"],
        DATA["Mastercard"]["intangibles_2025"],
    ]

    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E78")
        cell.alignment = Alignment(wrap_text=True, horizontal="center")
        ws.cell(row=2, column=col, value=values[col - 1])
        ws.column_dimensions[get_column_letter(col)].width = 18

    ratio_ws = wb.create_sheet("Ratios")
    ratio_ws.append(["Ratio", "Visa", "Mastercard", "Comment"])
    for cell in ratio_ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="0F766E")
        cell.alignment = Alignment(horizontal="center")

    for idx, (name, _group, visa_formula, ma_formula) in enumerate(RATIOS, start=2):
        ratio_ws.cell(row=idx, column=1, value=name)
        if visa_formula is None:
            ratio_ws.cell(row=idx, column=2, value="n/a")
            ratio_ws.cell(row=idx, column=3, value="n/a")
        else:
            ratio_ws.cell(row=idx, column=2, value=visa_formula.format(row=2))
            ratio_ws.cell(row=idx, column=3, value=ma_formula.format(row=2))
        ratio_ws.cell(row=idx, column=4, value=COMMENTS[name])

    for row_idx in [5, 6, 10, 11, 12, 19]:
        ratio_ws.cell(row=row_idx, column=2).number_format = "0.00%"
        ratio_ws.cell(row=row_idx, column=3).number_format = "0.00%"
    for row_idx in [2, 3, 4]:
        ratio_ws.cell(row=row_idx, column=2).number_format = "#,##0"
        ratio_ws.cell(row=row_idx, column=3).number_format = "#,##0"
    for row_idx in [7, 8, 9, 13, 14, 16, 17, 18]:
        ratio_ws.cell(row=row_idx, column=2).number_format = "0.00"
        ratio_ws.cell(row=row_idx, column=3).number_format = "0.00"

    for col in ["A", "B", "C", "D"]:
        ratio_ws.column_dimensions[col].width = 24 if col != "D" else 70

    analysis_ws = wb.create_sheet("Analysis")
    analysis_rows = [
        ("Main conclusion", "Visa looks stronger on financial quality; Mastercard looks stronger on efficiency."),
        ("Why Visa stands out", "Higher net margin, stronger interest coverage, and much lower leverage relative to equity."),
        ("Why Mastercard stands out", "Higher ROA, higher asset turnover, and higher fixed-asset turnover."),
        ("Critical caution", "Mastercard's ROE and price-to-book are heavily distorted by low accounting equity after large buybacks."),
        ("Asset structure", "Visa has a much larger goodwill and intangibles base, which affects asset composition and turnover interpretation."),
        ("Inventory ratio", "Inventory turnover is not meaningful here because both firms are service-based payment networks and do not report a meaningful inventory base."),
        ("Market data note", "Market prices used for P/E and Price-to-Book were checked on official investor quote pages on 2026-04-24."),
    ]
    analysis_ws.append(["Topic", "Takeaway"])
    for cell in analysis_ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="7C2D12")
    for row in analysis_rows:
        analysis_ws.append(row)
    analysis_ws.column_dimensions["A"].width = 24
    analysis_ws.column_dimensions["B"].width = 90

    source_ws = wb.create_sheet("Sources")
    source_ws.append(["Type", "Source"])
    for cell in source_ws[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="374151")
    sources = [
        ("Visa 10-K 2025", "https://www.sec.gov/Archives/edgar/data/1403161/000140316125000089/v-20250930.htm"),
        ("Mastercard 10-K 2025", "https://www.sec.gov/Archives/edgar/data/1141391/000114139126000013/ma-20251231.htm"),
        ("Visa IR annual report", "https://annualreport.visa.com/annual-meeting/default.aspx"),
        ("Mastercard IR SEC filings", "https://investor.mastercard.com/financials-and-sec-filings/sec-filings/sec-filings-details/default.aspx?FilingId=19135441"),
        ("Visa market quote reference page", "https://investor.visa.com/stock-information/quote-chart/"),
        ("Mastercard market quote reference page", "https://investor.mastercard.com/stock-information/stock-data-and-chart/default.aspx"),
        ("Market quote note", "P/E and Price to Book use market prices checked on 2026-04-24 against official investor quote pages."),
    ]
    for row in sources:
        source_ws.append(row)
    source_ws.column_dimensions["A"].width = 28
    source_ws.column_dimensions["B"].width = 120

    wb.save(OUT_DIR / "visa_vs_mastercard_analysis.xlsx")


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.7), Inches(12.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.color.rgb = RGBColor(15, 118, 110)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.6), Inches(1.0)).text_frame
    p2.text = subtitle
    p2.paragraphs[0].font.size = Pt(18)
    p2.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
    return slide


def add_bullet_slide(prs, title, bullets, accent=(31, 78, 121)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(252, 252, 249)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.45), Inches(12.2), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(*accent)
    banner.line.color.rgb = RGBColor(*accent)
    banner.text_frame.text = title
    banner.text_frame.paragraphs[0].font.size = Pt(24)
    banner.text_frame.paragraphs[0].font.bold = True
    banner.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.2), Inches(4.8)).text_frame
    body.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(17, 24, 39)
    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(11.5), Inches(0.4)).text_frame
    footer.text = "Sources: Visa 2025 Form 10-K, Mastercard 2025 Form 10-K, official IR pages, market quote checked on 2026-04-24."
    footer.paragraphs[0].font.size = Pt(9)
    footer.paragraphs[0].font.color.rgb = RGBColor(107, 114, 128)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Visa vs Mastercard",
        "Comparative financial analysis using FY2025 official filings and market prices checked on official investor quote pages on April 24, 2026.",
    )

    add_bullet_slide(
        prs,
        "1. Business Model",
        [
            "Both are payment network technology firms, not banks: they process and enable transactions rather than lend money.",
            "Visa's 2025 strategy emphasizes consumer payments, commercial and money movement solutions, and value-added services on top of VisaNet.",
            "Mastercard also combines network scale with value-added services, security, data, and cross-border capabilities.",
        ],
        accent=(15, 118, 110),
    )

    add_bullet_slide(
        prs,
        "2. Profitability And Risk",
        [
            "Visa has the stronger net margin: 50.1% vs 45.6%. It keeps more profit from each dollar of revenue.",
            "Visa is materially less leveraged: debt/assets 25.3% vs 35.1%, and long-term debt/equity 51.7% vs 235.9%.",
            "Visa also has stronger interest coverage: 40.74x vs 26.17x, so its debt burden looks safer.",
            "Mastercard's ROE is higher at 210.5%, but that is not pure operating superiority. It is heavily boosted by a very small equity base after buybacks.",
            "Conclusion on quality: Visa looks financially safer and cleaner.",
        ],
        accent=(124, 45, 18),
    )

    add_bullet_slide(
        prs,
        "3. Asset Structure",
        [
            "Net PP&E is higher at Visa: 4.236bn vs 2.303bn, but Mastercard still produces more revenue from its fixed-asset base.",
            "Visa carries far larger goodwill and intangible assets: about 47.525bn vs 14.942bn.",
            "That matters because a heavier intangible base can depress turnover ratios and changes how asset quality should be interpreted.",
            "Mastercard has lower absolute long-term debt than Visa, but looks riskier because its equity base is much thinner.",
        ],
        accent=(31, 78, 121),
    )

    add_bullet_slide(
        prs,
        "4. Efficiency",
        [
            "Mastercard is more efficient in using assets: ROA 29.3% vs 20.7%, asset turnover 0.64 vs 0.41, fixed-assets turnover 14.77 vs 9.93.",
            "Visa is stronger on receivables turnover: 14.07 vs 7.82, suggesting faster monetization of receivables.",
            "Inventory turnover is not meaningful for either company because they are service businesses with no meaningful inventory base.",
            "Conclusion on operations: Mastercard extracts more revenue and profit from a smaller asset base.",
        ],
        accent=(31, 78, 121),
    )

    add_bullet_slide(
        prs,
        "5. Market View And Final Judgment",
        [
            "The market values both firms at a similar earnings multiple: P/E about 30x for each based on current prices and FY2025 EPS.",
            "Price-to-book is radically higher for Mastercard: 57.48 vs 13.75. This mostly reflects very low accounting equity, not necessarily better intrinsic value.",
            "Dividend payout is higher at Visa: 23.1% vs 18.4%, so Visa returns a larger share of profit to shareholders.",
            "Overall conclusion: Mastercard is the more efficient operator, but Visa appears to perform better financially on a balanced basis because profitability quality is high and leverage is much lower.",
        ],
        accent=(17, 94, 89),
    )

    prs.save(OUT_DIR / "visa_vs_mastercard_presentation.pptx")


def visual_add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246, 244, 238)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(0.55), Inches(12.4), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.color.rgb = RGBColor(15, 118, 110)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.85), Inches(11.6), Inches(0.8)).text_frame
    sub.text = subtitle
    sub.paragraphs[0].font.size = Pt(16)
    sub.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)

    q = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(4.9), Inches(1.6)).text_frame
    q.text = "Question"
    q.paragraphs[0].font.size = Pt(16)
    q.paragraphs[0].font.bold = True
    q.paragraphs[0].font.color.rgb = RGBColor(15, 118, 110)
    p = q.add_paragraph()
    p.text = "Which company performs better financially?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 24, 39)

    for x, w, text, fill, color in [
        (0.7, 2.6, "FY2025 official filings", (221, 243, 239), (15, 118, 110)),
        (3.45, 2.85, "6 slides / 8-10 minutes", (254, 243, 199), (146, 64, 14)),
    ]:
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.85), Inches(w), Inches(0.65))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(*fill)
        tag.line.color.rgb = RGBColor(*fill)
        tag.text_frame.text = text
        tag.text_frame.paragraphs[0].font.size = Pt(14)
        tag.text_frame.paragraphs[0].font.bold = True
        tag.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)
    return slide


def visual_add_header(slide, title, accent):
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.4), Inches(12.4), Inches(0.85))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(*accent)
    banner.line.color.rgb = RGBColor(*accent)
    banner.text_frame.text = title
    banner.text_frame.paragraphs[0].font.size = Pt(24)
    banner.text_frame.paragraphs[0].font.bold = True
    banner.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def visual_add_footer(slide):
    footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.9), Inches(11.8), Inches(0.25)).text_frame
    footer.text = "Sources: Visa 2025 Form 10-K, Mastercard 2025 Form 10-K, official IR pages."
    footer.paragraphs[0].font.size = Pt(8)
    footer.paragraphs[0].font.color.rgb = RGBColor(107, 114, 128)


def visual_add_slide(prs, title, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(252, 252, 249)
    visual_add_header(slide, title, accent)
    visual_add_footer(slide)
    return slide


def visual_add_company_card(slide, x, y, w, h, title, subtitle, color):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*color)
    box.line.color.rgb = RGBColor(*color)
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(245, 245, 245)


def visual_add_message_box(slide, x, y, w, h, title, lines, fill):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill)
    box.line.color.rgb = RGBColor(*fill)
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(31, 41, 55)


def visual_add_kpi_chip(slide, x, y, w, title, value, fill, text_color=(17, 24, 39)):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.72))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(*fill)
    chip.line.color.rgb = RGBColor(*fill)
    tf = chip.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(*text_color)
    p = tf.add_paragraph()
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_color)


def visual_add_score_box(slide, x, y, w, h, title, score, subtitle, fill):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill)
    box.line.color.rgb = RGBColor(*fill)
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)
    p = tf.add_paragraph()
    p.text = score
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 24, 39)
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(55, 65, 81)


def visual_add_center_note(slide, x, y, w, text, fill, color=(17, 24, 39), size=18):
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.72))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(*fill)
    note.line.color.rgb = RGBColor(*fill)
    note.text_frame.text = text
    note.text_frame.paragraphs[0].font.size = Pt(size)
    note.text_frame.paragraphs[0].font.bold = True
    note.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)


def visual_add_bar(slide, x, y, width, label, visa_text, visa_ratio, ma_text, ma_ratio):
    label_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.35), Inches(0.3)).text_frame
    label_box.text = label
    label_box.paragraphs[0].font.size = Pt(15)
    label_box.paragraphs[0].font.bold = True
    label_box.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)

    for idx, (name, value, ratio, color) in enumerate(
        [
            ("Visa", visa_text, visa_ratio, (15, 118, 110)),
            ("Mastercard", ma_text, ma_ratio, (185, 28, 28)),
        ]
    ):
        row_y = y + 0.34 + idx * 0.38
        name_tf = slide.shapes.add_textbox(Inches(x), Inches(row_y), Inches(1.1), Inches(0.25)).text_frame
        name_tf.text = name
        name_tf.paragraphs[0].font.size = Pt(12)
        name_tf.paragraphs[0].font.bold = True
        name_tf.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)

        track = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.15), Inches(row_y + 0.02), Inches(width), Inches(0.16))
        track.fill.solid()
        track.fill.fore_color.rgb = RGBColor(229, 231, 235)
        track.line.color.rgb = RGBColor(229, 231, 235)

        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.15), Inches(row_y + 0.02), Inches(width * ratio), Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.color.rgb = RGBColor(*color)

        val = slide.shapes.add_textbox(Inches(x + 1.15 + width + 0.1), Inches(row_y - 0.03), Inches(1.25), Inches(0.25)).text_frame
        val.text = value
        val.paragraphs[0].font.size = Pt(12)
        val.paragraphs[0].font.bold = True
        val.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    visual_add_title_slide(
        prs,
        "Visa vs Mastercard",
        "FY2025 comparative analysis based on official filings, with market prices checked on official investor quote pages on April 24, 2026.",
    )

    slide = visual_add_slide(prs, "1. Same Industry, Different Financial Profile", (15, 118, 110))
    visual_add_company_card(slide, 0.8, 1.45, 2.8, 1.0, "Visa", "Scale + cleaner balance sheet", (15, 118, 110))
    visual_add_company_card(slide, 3.95, 1.45, 3.1, 1.0, "Mastercard", "Efficiency + thinner equity base", (185, 28, 28))
    visual_add_center_note(slide, 0.8, 2.8, 5.2, "Same business model: global payment networks, not banks", (221, 243, 239))
    visual_add_center_note(slide, 6.25, 2.8, 5.0, "The comparison is really about margins, leverage and asset use", (254, 243, 199), size=16)
    visual_add_kpi_chip(slide, 0.8, 4.0, 1.9, "Revenue", "$39.1bn", (230, 244, 241))
    visual_add_kpi_chip(slide, 2.95, 4.0, 1.9, "Net income", "$19.6bn", (230, 244, 241))
    visual_add_kpi_chip(slide, 5.1, 4.0, 1.9, "Revenue", "$32.0bn", (254, 226, 226))
    visual_add_kpi_chip(slide, 7.25, 4.0, 1.9, "Net income", "$14.6bn", (254, 226, 226))
    visual_add_score_box(slide, 9.65, 3.75, 2.1, 1.65, "Visual cue", "Visa", "bigger scale", (221, 243, 239))
    visual_add_center_note(slide, 1.2, 5.45, 10.5, "One sentence to say out loud: Visa is larger and safer; Mastercard is leaner and more aggressive.", (17, 24, 39), color=(255, 255, 255), size=18)

    slide = visual_add_slide(prs, "2. Profitability And Risk", (124, 45, 18))
    visual_add_bar(slide, 0.8, 1.45, 3.15, "Net margin", "50.1%", 1.0, "45.6%", 0.91)
    visual_add_bar(slide, 0.8, 2.55, 3.15, "Debt / assets", "25.3%", 0.72, "35.1%", 1.0)
    visual_add_bar(slide, 0.8, 3.65, 3.15, "Interest coverage", "40.7x", 1.0, "26.2x", 0.64)
    visual_add_score_box(slide, 6.2, 1.45, 2.35, 1.65, "Winner", "Visa", "better margin quality", (221, 243, 239))
    visual_add_score_box(slide, 8.85, 1.45, 2.35, 1.65, "Risk", "Visa", "safer leverage profile", (221, 243, 239))
    visual_add_bar(slide, 6.2, 3.45, 2.8, "LT debt / equity", "51.7%", 0.22, "235.9%", 1.0)
    visual_add_bar(slide, 9.25, 3.45, 2.15, "ROE", "52.1%", 0.25, "210.5%", 1.0)
    visual_add_center_note(slide, 1.0, 5.55, 10.9, "Key point: Mastercard's very high ROE looks impressive, but it is boosted by a much smaller equity base.", (254, 226, 226), size=17)

    slide = visual_add_slide(prs, "3. Asset Structure", (31, 78, 121))
    visual_add_bar(slide, 0.8, 1.45, 4.1, "Net PP&E", "4.24bn", 1.0, "2.30bn", 0.54)
    visual_add_bar(slide, 0.8, 2.75, 4.1, "Goodwill + intangibles", "47.53bn", 1.0, "14.94bn", 0.31)
    visual_add_bar(slide, 0.8, 4.05, 4.1, "Total LT debt", "19.60bn", 1.0, "18.25bn", 0.93)
    visual_add_score_box(slide, 6.15, 1.7, 2.45, 1.7, "Biggest gap", "Visa", "much heavier intangible base", (219, 234, 254))
    visual_add_score_box(slide, 8.95, 1.7, 2.45, 1.7, "Read carefully", "Mastercard", "less debt in dollars, more risky vs equity", (231, 240, 254))
    visual_add_center_note(slide, 6.15, 4.2, 5.25, "This slide explains why turnover ratios must be interpreted with the asset mix in mind.", (221, 243, 239), size=16)

    slide = visual_add_slide(prs, "4. Efficiency", (31, 78, 121))
    visual_add_bar(slide, 0.8, 1.45, 3.35, "ROA", "20.7%", 0.71, "29.3%", 1.0)
    visual_add_bar(slide, 0.8, 2.55, 3.35, "Asset turnover", "0.41", 0.64, "0.64", 1.0)
    visual_add_bar(slide, 0.8, 3.65, 3.35, "Fixed asset turnover", "9.93", 0.67, "14.77", 1.0)
    visual_add_bar(slide, 0.8, 4.75, 3.35, "Receivables turnover", "14.07", 1.0, "7.82", 0.56)
    visual_add_score_box(slide, 6.4, 1.55, 2.35, 1.65, "Efficiency", "Mastercard", "wins 3 of 4", (254, 226, 226))
    visual_add_score_box(slide, 9.0, 1.55, 2.35, 1.65, "Collection", "Visa", "faster receivables turn", (221, 243, 239))
    visual_add_center_note(slide, 6.1, 3.75, 5.65, "Inventory turnover is not meaningful here because both companies are service payment networks.", (220, 252, 231), size=15)

    slide = visual_add_slide(prs, "5. Market View And Final Judgment", (15, 118, 110))
    visual_add_bar(slide, 0.8, 1.35, 3.5, "P/E", "30.28x", 1.0, "30.41x", 1.0)
    visual_add_bar(slide, 0.8, 2.45, 3.5, "Price / book", "13.75x", 0.24, "57.48x", 1.0)
    visual_add_bar(slide, 0.8, 3.55, 3.5, "Dividend payout", "23.1%", 1.0, "18.4%", 0.80)
    visual_add_score_box(slide, 6.1, 1.25, 2.55, 1.8, "Overall winner", "Visa", "better financial quality", (221, 243, 239))
    visual_add_score_box(slide, 8.95, 1.25, 2.55, 1.8, "Operational edge", "Mastercard", "better asset efficiency", (254, 226, 226))
    visual_add_kpi_chip(slide, 6.1, 3.55, 2.4, "Visa wins on", "margin, leverage, payout", (230, 244, 241))
    visual_add_kpi_chip(slide, 8.95, 3.55, 2.4, "Mastercard wins on", "ROA and turnover", (254, 235, 235))
    verdict = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(5.75), Inches(10.9), Inches(0.7))
    verdict.fill.solid()
    verdict.fill.fore_color.rgb = RGBColor(17, 24, 39)
    verdict.line.color.rgb = RGBColor(17, 24, 39)
    verdict.text_frame.text = "Final call: Mastercard is more efficient operationally, but Visa performs better financially overall."
    verdict.text_frame.paragraphs[0].font.size = Pt(18)
    verdict.text_frame.paragraphs[0].font.bold = True
    verdict.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    prs.save(OUT_DIR / "visa_vs_mastercard_presentation.pptx")


class RGBColor:
    def __new__(cls, r, g, b):
        from pptx.dml.color import RGBColor as _RGBColor

        return _RGBColor(r, g, b)


if __name__ == "__main__":
    build_workbook()
    build_presentation()
